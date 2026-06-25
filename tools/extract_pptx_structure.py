#!/usr/bin/env python3
"""Extract PPTX slide text, embedded images, and slide screenshots to YAML.

Usage:
    tools/extract_pptx_structure.py courses/giao_duc_phap_luat/source-ppt/bai_01.pptx
    tools/extract_pptx_structure.py courses/giao_duc_phap_luat/source-ppt/bai_01.pptx --out courses/giao_duc_phap_luat/output/bai_01.yaml

PPTX structure is read with python-pptx. Slide screenshots are rendered by
converting the full PPTX to a temporary PDF with LibreOffice (`soffice`), then
rasterizing that PDF with Poppler (`pdftoppm`).
"""

from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as exc:  # pragma: no cover - exercised by CLI environment.
    raise SystemExit(
        "Missing dependency: python-pptx. Activate your environment first, "
        "for example: `conda activate ds`."
    ) from exc


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT_ROOT = ROOT / "pptx-analysis"


def markdown_text(text: str, bold: bool) -> str:
    if not bold or not text.strip():
        return text
    leading_len = len(text) - len(text.lstrip())
    trailing_len = len(text) - len(text.rstrip())
    leading = text[:leading_len]
    trailing = text[len(text) - trailing_len :] if trailing_len else ""
    core = text[leading_len : len(text) - trailing_len if trailing_len else len(text)]
    return f"{leading}**{core}**{trailing}"


def merge_segments(segments: list[tuple[str, bool]]) -> list[tuple[str, bool]]:
    merged: list[tuple[str, bool]] = []
    for text, bold in segments:
        text = text.replace("\t", "  ")
        if not text:
            continue
        if merged and merged[-1][1] == bold:
            prev_text, prev_bold = merged[-1]
            merged[-1] = (prev_text + text, prev_bold)
        else:
            merged.append((text, bold))
    return merged


def most_common_font_size(font_sizes: list[float]) -> float | None:
    if not font_sizes:
        return None
    counts: dict[float, int] = {}
    for size in font_sizes:
        counts[size] = counts.get(size, 0) + 1
    return sorted(counts.items(), key=lambda item: (-item[1], -item[0]))[0][0]


def clean_text_block(lines: list[str]) -> str:
    cleaned = [line.strip() for line in lines]
    while cleaned and cleaned[0] == "":
        cleaned.pop(0)
    while cleaned and cleaned[-1] == "":
        cleaned.pop()
    return "\n".join(cleaned)


def shape_order_key(shape: Any) -> tuple[int, int]:
    return (int(getattr(shape, "top", 0) or 0), int(getattr(shape, "left", 0) or 0))


def iter_shapes(shapes: Any) -> Any:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def extract_text_frame_lines(text_frame: Any) -> tuple[list[str], list[float]]:
    lines: list[str] = []
    font_sizes: list[float] = []
    for paragraph in text_frame.paragraphs:
        segments: list[tuple[str, bool]] = []
        for run in paragraph.runs:
            text = run.text or ""
            font = run.font
            if font.size is not None:
                font_sizes.append(font.size.pt)
            segments.append((text, bool(font.bold)))
        line = "".join(markdown_text(text, bold) for text, bold in merge_segments(segments))
        if line.strip():
            lines.append(line)
    return lines, font_sizes


def extract_table_lines(shape: Any) -> tuple[list[str], list[float]]:
    lines: list[str] = []
    font_sizes: list[float] = []
    if not getattr(shape, "has_table", False):
        return lines, font_sizes

    for row in shape.table.rows:
        cells: list[str] = []
        for cell in row.cells:
            cell_lines, cell_font_sizes = extract_text_frame_lines(cell.text_frame)
            font_sizes.extend(cell_font_sizes)
            cells.append(" ".join(line.strip() for line in cell_lines if line.strip()))
        row_text = " | ".join(cell for cell in cells if cell)
        if row_text:
            lines.append(row_text)
    return lines, font_sizes


def extract_text_blocks(slide: Any) -> list[dict[str, Any]]:
    raw_blocks: list[tuple[tuple[int, int], dict[str, Any]]] = []
    for shape in iter_shapes(slide.shapes):
        lines: list[str] = []
        font_sizes: list[float] = []
        if getattr(shape, "has_text_frame", False):
            lines, font_sizes = extract_text_frame_lines(shape.text_frame)
        elif getattr(shape, "has_table", False):
            lines, font_sizes = extract_table_lines(shape)

        if lines:
            text = clean_text_block(lines)
            if not text:
                continue
            block: dict[str, Any] = {
                "font_size_pt": most_common_font_size(font_sizes),
                "text": text,
            }
            raw_blocks.append((shape_order_key(shape), block))

    blocks = [block for _, block in sorted(raw_blocks, key=lambda item: item[0])]
    for index, block in enumerate(blocks, start=1):
        block["index"] = index
        ordered = {"index": block.pop("index")}
        if block.get("font_size_pt") is not None:
            ordered["font_size_pt"] = block["font_size_pt"]
        ordered["text"] = block["text"]
        blocks[index - 1] = ordered
    return blocks


def extract_images(slide: Any, image_dir: Path, slide_number: int) -> list[dict[str, Any]]:
    raw_images: list[tuple[tuple[int, int], dict[str, Any]]] = []
    pic_index = 0

    for shape in iter_shapes(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        pic_index += 1
        image = shape.image
        blob = image.blob
        suffix = image_suffix(shape, blob)
        out_name = f"slide_{slide_number:03d}_image_{pic_index:02d}{suffix}"
        out_path = image_dir / out_name
        out_path.write_bytes(blob)
        image_data: dict[str, Any] = {
            "file": relpath_for_yaml(out_path, image_dir.parent),
        }
        size = image_size(blob)
        if size:
            image_data["width_px"] = size[0]
            image_data["height_px"] = size[1]
        raw_images.append((shape_order_key(shape), image_data))

    images = [image for _, image in sorted(raw_images, key=lambda item: item[0])]
    for index, image in enumerate(images, start=1):
        image["index"] = index
        ordered = {"index": image.pop("index"), "file": image["file"]}
        if "width_px" in image:
            ordered["width_px"] = image["width_px"]
        if "height_px" in image:
            ordered["height_px"] = image["height_px"]
        images[index - 1] = ordered
    return images


def image_suffix(shape: Any, data: bytes) -> str:
    suffix = image_suffix_from_bytes(data)
    if suffix:
        return suffix

    part = image_part(shape)
    if part is not None:
        partname = getattr(part, "partname", None)
        if partname is not None:
            suffix = Path(str(partname)).suffix
            if suffix:
                return suffix

        content_type = getattr(part, "content_type", "")
        suffix = {
            "image/bmp": ".bmp",
            "image/gif": ".gif",
            "image/jpeg": ".jpg",
            "image/png": ".png",
            "image/tiff": ".tiff",
            "image/x-emf": ".emf",
            "image/x-wmf": ".wmf",
            "image/svg+xml": ".svg",
        }.get(content_type)
        if suffix:
            return suffix

    return ".bin"


def image_part(shape: Any) -> Any | None:
    try:
        rel_id = shape._element.blip_rId
        return shape.part.related_part(rel_id)
    except Exception:
        return None


def image_suffix_from_bytes(data: bytes) -> str | None:
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return ".png"
    if data.startswith(b"\xff\xd8"):
        return ".jpg"
    if data.startswith(b"GIF87a") or data.startswith(b"GIF89a"):
        return ".gif"
    if data.startswith(b"BM"):
        return ".bmp"
    if data.startswith(b"II*\x00") or data.startswith(b"MM\x00*"):
        return ".tiff"
    if data.startswith(b"\x01\x00\x00\x00"):
        return ".emf"
    if data.startswith(b"\xd7\xcd\xc6\x9a"):
        return ".wmf"
    return None


def image_size(data: bytes) -> tuple[int, int] | None:
    if data.startswith(b"\x89PNG\r\n\x1a\n") and len(data) >= 24:
        return (int.from_bytes(data[16:20], "big"), int.from_bytes(data[20:24], "big"))
    if (data.startswith(b"GIF87a") or data.startswith(b"GIF89a")) and len(data) >= 10:
        return (int.from_bytes(data[6:8], "little"), int.from_bytes(data[8:10], "little"))
    if data.startswith(b"\xff\xd8"):
        index = 2
        while index + 9 < len(data):
            if data[index] != 0xFF:
                index += 1
                continue
            marker = data[index + 1]
            index += 2
            if marker in (0xD8, 0xD9):
                continue
            if index + 2 > len(data):
                break
            length = int.from_bytes(data[index : index + 2], "big")
            if 0xC0 <= marker <= 0xC3 and index + 7 < len(data):
                height = int.from_bytes(data[index + 3 : index + 5], "big")
                width = int.from_bytes(data[index + 5 : index + 7], "big")
                return (width, height)
            index += length
    return None


def cleanup_previous_screenshots(screenshot_dir: Path) -> None:
    for pattern in ("slide_*.png", "slide_*.jpg", "slide-*.png", "slide-*.jpg"):
        for path in screenshot_dir.glob(pattern):
            path.unlink()


def render_slide_screenshots(
    pptx_path: Path,
    screenshot_dir: Path,
    dpi: int,
    image_format: str,
) -> list[Path]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("Cannot render screenshots: `soffice` or `libreoffice` was not found.")

    screenshot_dir.mkdir(parents=True, exist_ok=True)
    cleanup_previous_screenshots(screenshot_dir)
    with tempfile.TemporaryDirectory(prefix="pptx-render-") as tmp:
        tmp_path = Path(tmp)
        pdf_path = convert_pptx_to_pdf(pptx_path, tmp_path, soffice)
        render_pdf_pages(pdf_path, screenshot_dir, dpi, image_format)

    suffix = "jpg" if image_format == "jpg" else "png"
    rendered = sorted(screenshot_dir.glob(f"slide-*.{suffix}"))
    for index, path in enumerate(rendered, start=1):
        normalized = screenshot_dir / f"slide_{index:03d}.{suffix}"
        if path != normalized:
            path.replace(normalized)
    return sorted(screenshot_dir.glob(f"slide_*.{suffix}"))


def convert_pptx_to_pdf(pptx_path: Path, tmp_path: Path, soffice: str) -> Path:
    profile_dir = tmp_path / "lo-profile"
    pdf_dir = tmp_path / "pdf"
    pdf_dir.mkdir()
    cmd = [
        soffice,
        f"-env:UserInstallation=file://{profile_dir}",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pdf_dir),
        str(pptx_path),
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    pdfs = sorted(pdf_dir.glob("*.pdf"))
    if not pdfs:
        raise RuntimeError("LibreOffice did not produce a PDF for screenshot rendering.")
    return pdfs[0]


def render_pdf_pages(pdf_path: Path, screenshot_dir: Path, dpi: int, image_format: str) -> None:
    errors: list[str] = []

    pdftoppm = shutil.which("pdftoppm")
    if pdftoppm:
        try:
            render_pdf_pages_with_pdftoppm(pdftoppm, pdf_path, screenshot_dir, dpi, image_format)
            return
        except subprocess.CalledProcessError as exc:
            errors.append(command_error("pdftoppm", exc))

    gs = shutil.which("gs")
    if gs:
        try:
            render_pdf_pages_with_ghostscript(gs, pdf_path, screenshot_dir, dpi, image_format)
            return
        except subprocess.CalledProcessError as exc:
            errors.append(command_error("gs", exc))

    if not errors:
        raise RuntimeError("Cannot render screenshots: neither `pdftoppm` nor `gs` was found.")
    raise RuntimeError("Cannot render PDF screenshots. " + " | ".join(errors))


def render_pdf_pages_with_pdftoppm(
    pdftoppm: str,
    pdf_path: Path,
    screenshot_dir: Path,
    dpi: int,
    image_format: str,
) -> None:
    prefix = screenshot_dir / "slide"
    if image_format == "jpg":
        cmd = [
            pdftoppm,
            "-jpeg",
            "-jpegopt",
            "quality=85,progressive=y,optimize=y",
            "-r",
            str(dpi),
            str(pdf_path),
            str(prefix),
        ]
    else:
        cmd = [pdftoppm, "-png", "-r", str(dpi), str(pdf_path), str(prefix)]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)


def render_pdf_pages_with_ghostscript(
    gs: str,
    pdf_path: Path,
    screenshot_dir: Path,
    dpi: int,
    image_format: str,
) -> None:
    if image_format == "jpg":
        device = "jpeg"
        output = screenshot_dir / "slide-%03d.jpg"
        quality_args = ["-dJPEGQ=85"]
    else:
        device = "png16m"
        output = screenshot_dir / "slide-%03d.png"
        quality_args = []

    cmd = [
        gs,
        "-dSAFER",
        "-dBATCH",
        "-dNOPAUSE",
        f"-sDEVICE={device}",
        f"-r{dpi}",
        *quality_args,
        f"-sOutputFile={output}",
        str(pdf_path),
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)


def command_error(tool: str, exc: subprocess.CalledProcessError) -> str:
    stderr = (exc.stderr or "").strip().splitlines()
    detail = stderr[-1] if stderr else str(exc)
    return f"{tool} failed: {detail}"


def yaml_scalar(value: str) -> str:
    if value == "":
        return "''"
    if "\t" in value:
        escaped = (
            value.replace("\\", "\\\\")
            .replace('"', '\\"')
            .replace("\t", "\\t")
            .replace("\r", "\\r")
            .replace("\n", "\\n")
        )
        return f'"{escaped}"'
    if "\n" in value:
        lines = value.splitlines()
        if value.endswith("\n"):
            lines.append("")
        return "|\n" + "\n".join(f"  {line}" for line in lines)
    escaped = value.replace("\\", "\\\\").replace('"', '\\"')
    return f'"{escaped}"'


def to_yaml(value: Any, indent: int = 0) -> str:
    space = " " * indent
    if value is None:
        return "null"
    if isinstance(value, bool):
        return "true" if value else "false"
    if isinstance(value, (int, float)):
        return str(value)
    if isinstance(value, str):
        return yaml_scalar(value)
    if isinstance(value, list):
        if not value:
            return "[]"
        lines: list[str] = []
        for item in value:
            if isinstance(item, (dict, list)):
                lines.append(f"{space}- {to_yaml(item, indent + 2).lstrip()}")
            else:
                lines.append(f"{space}- {to_yaml(item, indent + 2)}")
        return "\n".join(lines)
    if isinstance(value, dict):
        if not value:
            return "{}"
        lines = []
        for key, item in value.items():
            if isinstance(item, str) and "\n" in item:
                block = yaml_scalar(item).splitlines()
                lines.append(f"{space}{key}: {block[0]}")
                lines.extend(f"{space}{line}" for line in block[1:])
            elif isinstance(item, (dict, list)) and item:
                lines.append(f"{space}{key}:")
                lines.append(to_yaml(item, indent + 2))
            else:
                lines.append(f"{space}{key}: {to_yaml(item, indent + 2)}")
        return "\n".join(lines)
    return yaml_scalar(str(value))


def relpath_for_yaml(path: Path, base_dir: Path) -> str:
    return path.resolve().relative_to(base_dir.resolve()).as_posix()


def dot_relpath_for_yaml(path: Path, base_dir: Path) -> str:
    return f"./{relpath_for_yaml(path, base_dir)}"


def analyze_pptx(
    pptx_path: Path,
    output_yaml: Path,
    render_screenshots: bool,
    screenshot_dpi: int,
    screenshot_format: str,
) -> dict[str, Any]:
    out_dir = output_yaml.parent
    image_dir = out_dir / "images"
    screenshot_dir = out_dir / "screenshots"
    image_dir.mkdir(parents=True, exist_ok=True)
    screenshot_dir.mkdir(parents=True, exist_ok=True)

    screenshots: list[Path] = []
    screenshot_error: str | None = None
    if render_screenshots:
        try:
            screenshots = render_slide_screenshots(
                pptx_path=pptx_path,
                screenshot_dir=screenshot_dir,
                dpi=screenshot_dpi,
                image_format=screenshot_format,
            )
        except Exception as exc:  # Keep extraction useful even without renderer.
            screenshot_error = str(exc)

    slides: list[dict[str, Any]] = []
    presentation = Presentation(str(pptx_path))
    for slide_number, slide in enumerate(presentation.slides, start=1):
        screenshot_file = None
        if slide_number <= len(screenshots):
            screenshot_file = dot_relpath_for_yaml(screenshots[slide_number - 1], out_dir)

        slide_data: dict[str, Any] = {
            "number": slide_number,
            "screenshot": screenshot_file,
            "texts": extract_text_blocks(slide),
            "images": extract_images(
                slide=slide,
                image_dir=image_dir,
                slide_number=slide_number,
            ),
        }
        slides.append(slide_data)

    result: dict[str, Any] = {
        "slide_count": len(slides),
        "slides": slides,
    }
    if screenshot_error:
        result["screenshot_error"] = screenshot_error
    return result


def default_output_yaml(pptx_path: Path) -> Path:
    parts = pptx_path.resolve().parts
    if "courses" in parts:
        course_index = parts.index("courses") + 1
        if course_index < len(parts):
            course_root = Path(*parts[: course_index + 1])
            return course_root / "pptx-analysis" / pptx_path.stem / f"{pptx_path.stem}.yaml"

    return DEFAULT_OUTPUT_ROOT / pptx_path.stem / f"{pptx_path.stem}.yaml"


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Analyze a PPTX into YAML with formatted text, embedded images, and slide screenshots."
    )
    parser.add_argument("pptx", type=Path, help="Source .pptx file, e.g. courses/giao_duc_phap_luat/source-ppt/bai_01.pptx")
    parser.add_argument(
        "-o",
        "--out",
        type=Path,
        help="Output YAML path. Default: courses/giao_duc_phap_luat/pptx-analysis/<deck>/<deck>.yaml",
    )
    parser.add_argument(
        "--no-screenshots",
        action="store_true",
        help="Skip LibreOffice PDF conversion and Poppler screenshot rendering.",
    )
    parser.add_argument(
        "--screenshot-dpi",
        type=int,
        default=120,
        help="Rasterization DPI for slide screenshots after PDF conversion. Default: 120.",
    )
    parser.add_argument(
        "--screenshot-format",
        choices=("jpg", "png"),
        default="jpg",
        help="Screenshot image format after PDF conversion. Default: jpg for smaller files.",
    )
    args = parser.parse_args()

    pptx_path = args.pptx.resolve()
    if not pptx_path.exists():
        print(f"Missing PPTX file: {pptx_path}", file=sys.stderr)
        return 1
    if pptx_path.suffix.lower() != ".pptx":
        print(f"Expected a .pptx file: {pptx_path}", file=sys.stderr)
        return 1
    if args.screenshot_dpi < 36:
        print("--screenshot-dpi must be at least 36.", file=sys.stderr)
        return 1

    output_yaml = (args.out.resolve() if args.out else default_output_yaml(pptx_path).resolve())
    output_yaml.parent.mkdir(parents=True, exist_ok=True)

    analysis = analyze_pptx(
        pptx_path=pptx_path,
        output_yaml=output_yaml,
        render_screenshots=not args.no_screenshots,
        screenshot_dpi=args.screenshot_dpi,
        screenshot_format=args.screenshot_format,
    )
    output_yaml.write_text(to_yaml(analysis) + "\n", encoding="utf-8")
    print(output_yaml)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
